# -*- coding: utf-8 -*-
# =============================================================================
# Ag-ppt-create - AI-powered PPTX generation pipeline
# https://github.com/aktsmm/Ag-ppt-create
# 
# Copyright (c) aktsmm. Licensed under CC BY-NC-SA 4.0.
# DO NOT MODIFY THIS HEADER BLOCK.
# =============================================================================
"""
Create Japanese summary PPTX from content JSON.
Usage: python create_ja_pptx.py <content.json> <output.pptx> [--no-signature]

Supports image embedding:
  - "image": {"path": "images/foo.png", "position": "right", "width_percent": 45}
  - "image": {"url": "https://...", "position": "bottom", "height_percent": 50}
  - position: "right" | "bottom" | "full"
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import json
import sys
import os
import urllib.request
import tempfile
import subprocess
from pathlib import Path

# Repository signature
REPO_URL = "https://github.com/aktsmm/Ag-ppt-create"


def get_repo_info() -> str:
    """Get repository URL from git remote (fallback to default)."""
    try:
        result = subprocess.run(
            ['git', 'remote', 'get-url', 'origin'],
            capture_output=True, text=True, timeout=5
        )
        if result.returncode == 0:
            return result.stdout.strip()
    except Exception:
        pass
    return REPO_URL


# Parse arguments
add_signature = True
if '--no-signature' in sys.argv:
    add_signature = False
    sys.argv.remove('--no-signature')

if len(sys.argv) >= 3:
    input_json = sys.argv[1]
    output_pptx = sys.argv[2]
else:
    input_json = 'output_manifest/20251211_brk252_ja_summary_content.json'
    output_pptx = 'output_ppt/20251211_brk252_ja_summary_v2.pptx'

# Load content
with open(input_json, 'r', encoding='utf-8') as f:
    data = json.load(f)


def validate_and_fix_content(slides_data: list) -> list:
    """
    Validate content and auto-fix common issues.
    - closing type with items â†’ convert to content type
    """
    fixed_count = 0
    
    for i, slide in enumerate(slides_data):
        slide_type = slide.get('type', 'content')
        items = slide.get('items', slide.get('content', slide.get('content_ja', [])))
        title = slide.get('title', slide.get('title_ja', f'Slide {i + 1}'))
        
        # Check: closing type should not have multiple items
        if slide_type == 'closing' and items and len(items) > 1:
            print(f"  âš ï¸  Warning: Slide {i + 1} '{title[:30]}...' has type='closing' with {len(items)} items")
            print(f"      â†’ Auto-converting to type='content' (closing is for short endings only)")
            slide['type'] = 'content'
            fixed_count += 1
    
    if fixed_count > 0:
        print(f"  âœ… Auto-fixed {fixed_count} slide(s)")
    
    return slides_data


# Validate and fix content
data['slides'] = validate_and_fix_content(data['slides'])

# Create presentation
prs = Presentation()
prs.slide_width = Inches(13.333)  # 16:9
prs.slide_height = Inches(7.5)

# Colors
PURPLE = RGBColor(0x5B, 0x5F, 0xC7)  # Microsoft Purple
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)

# Get base directory for resolving relative image paths
BASE_DIR = Path(input_json).parent.parent  # Assume content.json is in output_manifest/


def resolve_image_path(image_config: dict) -> str:
    """
    Resolve image path from config. Supports local path or URL.
    Returns the resolved local file path (downloads URL if needed).
    """
    if 'url' in image_config:
        # Download from URL to temp file
        url = image_config['url']
        try:
            suffix = Path(url.split('?')[0]).suffix or '.png'
            with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as tmp:
                print(f"  ğŸ“¥ Downloading image: {url[:60]}...")
                urllib.request.urlretrieve(url, tmp.name)
                return tmp.name
        except Exception as e:
            print(f"  âš ï¸  Failed to download image: {e}")
            return None
    elif 'path' in image_config:
        # Resolve local path
        path = image_config['path']
        # Try relative to BASE_DIR first, then as absolute
        candidates = [
            BASE_DIR / path,
            Path(path),
            Path(input_json).parent / path,
        ]
        for candidate in candidates:
            if candidate.exists():
                return str(candidate)
        print(f"  âš ï¸  Image not found: {path}")
        return None
    return None


def add_image_to_slide(slide, image_config: dict, content_area: dict):
    """
    Add image to slide based on configuration.
    
    Args:
        slide: pptx slide object
        image_config: {"path": ..., "position": "right|bottom|full", "width_percent": 45}
        content_area: {"left": Inches, "top": Inches, "width": Inches, "height": Inches}
    
    Returns:
        Updated content_area for text placement
    """
    image_path = resolve_image_path(image_config)
    if not image_path:
        return content_area
    
    position = image_config.get('position', 'right')
    width_pct = image_config.get('width_percent', 45)
    height_pct = image_config.get('height_percent', 50)
    
    slide_width = float(prs.slide_width)
    slide_height = float(prs.slide_height)
    
    try:
        if position == 'full':
            # Full slide image (centered, maintaining aspect ratio)
            img_left = Inches(0.5)
            img_top = content_area['top']
            img_width = Inches(12.333)
            slide.shapes.add_picture(image_path, img_left, img_top, width=img_width)
            # No text area when full image
            return None
            
        elif position == 'right':
            # Image on right, text on left
            img_width = Inches(13.333 * width_pct / 100)
            img_left = Inches(13.333 - 0.5) - img_width
            img_top = content_area['top']
            slide.shapes.add_picture(image_path, img_left, img_top, width=img_width)
            # Update content area to left side
            return {
                'left': content_area['left'],
                'top': content_area['top'],
                'width': Inches(13.333 * (100 - width_pct - 5) / 100),
                'height': content_area['height']
            }
            
        elif position == 'bottom':
            # Image at bottom, text on top
            text_height = Inches(7.5 * (100 - height_pct - 10) / 100)
            img_height = Inches(7.5 * height_pct / 100)
            img_top = Inches(7.5 - 0.5) - img_height
            img_left = Inches(0.5)
            slide.shapes.add_picture(image_path, img_left, img_top, height=img_height)
            # Update content area to top
            return {
                'left': content_area['left'],
                'top': content_area['top'],
                'width': content_area['width'],
                'height': text_height
            }
            
    except Exception as e:
        print(f"  âš ï¸  Failed to add image: {e}")
    
    return content_area


def add_title_slide(prs, title, subtitle=''):
    """Add a title slide with purple background."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = PURPLE
    bg.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(2))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(24)
        p2.font.color.rgb = WHITE
        p2.alignment = PP_ALIGN.CENTER
    
    return slide


def add_content_slide(prs, title, items, slide_type='feature', image_config=None):
    """
    Add a content slide with title bar, bullet points, and optional image.
    
    Args:
        prs: Presentation object
        title: Slide title
        items: List of bullet points
        slide_type: Type of slide (for styling)
        image_config: Optional image configuration dict
    """
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Title bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    bar.fill.solid()
    bar.fill.fore_color.rgb = PURPLE
    bar.line.fill.background()
    
    # Title text
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # Default content area
    content_area = {
        'left': Inches(0.5),
        'top': Inches(1.5),
        'width': Inches(12.333),
        'height': Inches(5.5)
    }
    
    # Add image if configured (adjusts content_area)
    if image_config:
        content_area = add_image_to_slide(slide, image_config, content_area)
    
    # Content (skip if image takes full slide)
    if items and content_area:
        content_box = slide.shapes.add_textbox(
            content_area['left'],
            content_area['top'],
            content_area['width'],
            content_area['height']
        )
        tf = content_box.text_frame
        tf.word_wrap = True
        
        for i, item in enumerate(items):
            if i > 0:
                p = tf.add_paragraph()
            else:
                p = tf.paragraphs[0]
            
            if isinstance(item, dict):
                if 'stat' in item:
                    desc = item.get('description', item.get('text', ''))
                    p.text = f"â€¢ {item.get('stat', '')} - {desc}"
                elif 'label' in item:
                    p.text = f"â€¢ {item['label']} {item.get('text', '')}"
                else:
                    p.text = f"â€¢ {item.get('text', str(item))}"
            else:
                p.text = f"â€¢ {item}"
            
            p.font.size = Pt(24)
            p.font.color.rgb = DARK_GRAY
            p.space_after = Pt(12)
    
    return slide


def add_section_slide(prs, title, subtitle=''):
    """Add a section divider slide with purple background."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = PURPLE
    bg.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.8), Inches(12.333), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(28)
        p2.font.color.rgb = WHITE
        p2.alignment = PP_ALIGN.CENTER
    
    return slide


def add_two_column_slide(prs, title, left_title, left_items, right_title, right_items):
    """Add a two-column comparison slide."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Title bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    bar.fill.solid()
    bar.fill.fore_color.rgb = PURPLE
    bar.line.fill.background()
    
    # Title text
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # Left column title
    if left_title:
        left_title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(5.8), Inches(0.6))
        tf = left_title_box.text_frame
        p = tf.paragraphs[0]
        p.text = left_title
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = PURPLE
    
    # Left column content
    if left_items:
        left_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.2), Inches(5.8), Inches(4.8))
        tf = left_box.text_frame
        tf.word_wrap = True
        for i, item in enumerate(left_items):
            if i > 0:
                p = tf.add_paragraph()
            else:
                p = tf.paragraphs[0]
            text = item.get('text', item) if isinstance(item, dict) else item
            p.text = f"â€¢ {text}"
            p.font.size = Pt(20)
            p.font.color.rgb = DARK_GRAY
            p.space_after = Pt(10)
    
    # Right column title
    if right_title:
        right_title_box = slide.shapes.add_textbox(Inches(7), Inches(1.5), Inches(5.8), Inches(0.6))
        tf = right_title_box.text_frame
        p = tf.paragraphs[0]
        p.text = right_title
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = PURPLE
    
    # Right column content
    if right_items:
        right_box = slide.shapes.add_textbox(Inches(7), Inches(2.2), Inches(5.8), Inches(4.8))
        tf = right_box.text_frame
        tf.word_wrap = True
        for i, item in enumerate(right_items):
            if i > 0:
                p = tf.add_paragraph()
            else:
                p = tf.paragraphs[0]
            text = item.get('text', item) if isinstance(item, dict) else item
            p.text = f"â€¢ {text}"
            p.font.size = Pt(20)
            p.font.color.rgb = DARK_GRAY
            p.space_after = Pt(10)
    
    return slide


# Create slides
slides_created = []
for slide_data in data['slides']:
    slide_type = slide_data.get('type', 'feature')
    title = slide_data.get('title_ja', '') or slide_data.get('title', '')
    subtitle = slide_data.get('subtitle_ja', '') or slide_data.get('subtitle', '')
    content = slide_data.get('content_ja', []) or slide_data.get('content', []) or slide_data.get('items', [])
    image_config = slide_data.get('image')  # Optional image configuration
    notes = slide_data.get('notes', '')  # Speaker notes
    
    if slide_type in ['title', 'closing']:
        slide = add_title_slide(prs, title, subtitle)
    elif slide_type in ['section', 'section_title']:
        slide = add_section_slide(prs, title, subtitle)
    elif slide_type == 'two_column':
        left_title = slide_data.get('left_title', '')
        left_items = slide_data.get('left_items', [])
        right_title = slide_data.get('right_title', '')
        right_items = slide_data.get('right_items', [])
        slide = add_two_column_slide(prs, title, left_title, left_items, right_title, right_items)
    else:
        slide = add_content_slide(prs, title, content, slide_type, image_config)
    
    # Add speaker notes if provided
    if notes and slide:
        try:
            notes_slide = slide.notes_slide
            notes_tf = notes_slide.notes_text_frame
            notes_tf.text = notes
        except Exception:
            pass
    
    slides_created.append(slide)

# Add signature to first and last slides
if add_signature and len(slides_created) > 0:
    repo_url = get_repo_info()
    sig_first = f"ğŸ“Œ Generated by: {repo_url}"
    sig_last = f"---\nğŸ”§ This presentation was created using Ag-ppt-create\n{repo_url}"
    
    # First slide
    try:
        notes_slide = slides_created[0].notes_slide
        notes_tf = notes_slide.notes_text_frame
        existing = notes_tf.text or ""
        if sig_first not in existing:
            notes_tf.text = f"{sig_first}\n\n{existing}" if existing else sig_first
    except Exception:
        pass
    
    # Last slide (if different from first)
    if len(slides_created) > 1:
        try:
            notes_slide = slides_created[-1].notes_slide
            notes_tf = notes_slide.notes_text_frame
            existing = notes_tf.text or ""
            if sig_last not in existing:
                notes_tf.text = f"{existing}\n\n{sig_last}" if existing else sig_last
        except Exception:
            pass
    
    print(f'ğŸ“ Signature added to speaker notes')

# Save
prs.save(output_pptx)
print(f'ä¿å­˜å®Œäº†: {output_pptx}')
print(f'ã‚¹ãƒ©ã‚¤ãƒ‰æ•°: {len(prs.slides)}')
