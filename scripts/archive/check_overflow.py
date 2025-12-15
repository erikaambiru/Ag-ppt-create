# -*- coding: utf-8 -*-
# =============================================================================
# Ag-ppt-create - AI-powered PPTX generation pipeline
# https://github.com/aktsmm/Ag-ppt-create
# 
# Copyright (c) aktsmm. Licensed under CC BY-NC-SA 4.0.
# DO NOT MODIFY THIS HEADER BLOCK.
# =============================================================================
"""Check text overflow in PPTX file."""
from pptx import Presentation
from pptx.util import Pt
import sys

def check_overflow(pptx_path):
    """Check for text overflow issues in PPTX."""
    prs = Presentation(pptx_path)
    
    print(f"=== PPTX Quality Check ===")
    print(f"File: {pptx_path}")
    print(f"Total slides: {len(prs.slides)}\n")
    
    issues = []
    
    for slide_idx, slide in enumerate(prs.slides):
        for shape_idx, shape in enumerate(slide.shapes):
            if not hasattr(shape, "text_frame") or not shape.has_text_frame:
                continue
                
            text_frame = shape.text_frame
            shape_id = f"Slide {slide_idx}, Shape {shape_idx}"
            
            # Check text length
            text_content = text_frame.text
            if len(text_content) > 500:
                issues.append(f"⚠️  {shape_id}: Text too long ({len(text_content)} chars)")
            
            # Check paragraph count
            para_count = len(text_frame.paragraphs)
            if para_count > 15:
                issues.append(f"⚠️  {shape_id}: Too many paragraphs ({para_count})")
            
            # Preview content
            preview = text_content[:80].replace('\n', ' ')
            print(f"{shape_id}: {preview}...")
    
    print("\n=== Issues Found ===")
    if issues:
        for issue in issues:
            print(issue)
        return False
    else:
        print("✅ No overflow or excessive text detected")
        return True

if __name__ == "__main__":
    pptx_file = sys.argv[1] if len(sys.argv) > 1 else "output_ppt/20251211_microsoft_intro_report.pptx"
    success = check_overflow(pptx_file)
    sys.exit(0 if success else 1)
