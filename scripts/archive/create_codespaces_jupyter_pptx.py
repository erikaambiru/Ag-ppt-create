# -*- coding: utf-8 -*-
"""
Create Japanese PPTX from content JSON for Codespaces Jupyter LT.
Usage: python scripts/create_codespaces_jupyter_pptx.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import json
import sys

# Load content
content_path = 'output_manifest/20251212_codespaces_jupyter_lt_content.json'
output_path = 'output_ppt/20251212_codespaces_jupyter_lt_template.pptx'

with open(content_path, 'r', encoding='utf-8') as f:
    data = json.load(f)

# Create presentation (16:9)
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Colors - GitHub theme
GITHUB_DARK = RGBColor(0x24, 0x29, 0x2e)
GITHUB_BLUE = RGBColor(0x03, 0x66, 0xd6)
JUPYTER_ORANGE = RGBColor(0xf3, 0x76, 0x26)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF6, 0xF8, 0xFA)


def add_title_slide(prs, title, subtitle=''):
    """Add a title slide with gradient-like background."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = GITHUB_DARK
    bg.line.fill.background()
    
    # Accent bar
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(6.5), prs.slide_width, Inches(1))
    accent.fill.solid()
    accent.fill.fore_color.rgb = GITHUB_BLUE
    accent.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.0), Inches(12.333), Inches(2.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(24)
        p2.font.color.rgb = RGBColor(0x79, 0xb8, 0xff)  # Light blue
        p2.alignment = PP_ALIGN.CENTER
        p2.space_before = Pt(20)
    
    return slide


def add_agenda_slide(prs, title, items):
    """Add an agenda slide."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Title bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    bar.fill.solid()
    bar.fill.fore_color.rgb = GITHUB_BLUE
    bar.line.fill.background()
    
    # Title text
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"📋 {title}"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # Content box
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(11.5), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    for i, item in enumerate(items):
        if i > 0:
            p = tf.add_paragraph()
        else:
            p = tf.paragraphs[0]
        p.text = f"{i+1}. {item}"
        p.font.size = Pt(26)
        p.font.color.rgb = DARK_GRAY
        p.space_before = Pt(16)
        p.space_after = Pt(8)
    
    return slide


def add_content_slide(prs, title, items, icon='📌'):
    """Add a content slide with bullet points."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Title bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    bar.fill.solid()
    bar.fill.fore_color.rgb = GITHUB_BLUE
    bar.line.fill.background()
    
    # Title text
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"{icon} {title}"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # Content box
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(11.5), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    for i, item in enumerate(items):
        if i > 0:
            p = tf.add_paragraph()
        else:
            p = tf.paragraphs[0]
        p.text = f"• {item}"
        p.font.size = Pt(22)
        p.font.color.rgb = DARK_GRAY
        p.space_before = Pt(12)
        p.space_after = Pt(6)
    
    return slide


def add_summary_slide(prs, title, items):
    """Add a summary slide with dark background."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = GITHUB_DARK
    bg.line.fill.background()
    
    # Accent bar at top
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.1))
    accent.fill.solid()
    accent.fill.fore_color.rgb = GITHUB_BLUE
    accent.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.333), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"📝 {title}"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # Content box
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.5), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    for i, item in enumerate(items):
        if i > 0:
            p = tf.add_paragraph()
        else:
            p = tf.paragraphs[0]
        p.text = f"✅ {item}"
        p.font.size = Pt(22)
        p.font.color.rgb = RGBColor(0xe1, 0xe4, 0xe8)  # Light gray
        p.space_before = Pt(10)
        p.space_after = Pt(6)
    
    return slide


# Build slides
icons = {
    "Jupyter Notebook とは": "📓",
    "機械学習環境の選択肢": "☁️",
    "GitHub Codespaces とは": "🐙",
    "セットアップ手順（ワンクリック！）": "🚀",
    "ハンズオン：実際に動かしてみよう": "🔧",
}

for slide_data in data['slides']:
    slide_type = slide_data['type']
    title = slide_data['title']
    
    if slide_type == 'title':
        add_title_slide(prs, title, slide_data.get('subtitle', ''))
    elif slide_type == 'agenda':
        add_agenda_slide(prs, title, slide_data['items'])
    elif slide_type == 'summary':
        add_summary_slide(prs, title, slide_data['items'])
    else:  # content
        icon = icons.get(title, '📌')
        add_content_slide(prs, title, slide_data['items'], icon)

# Save
prs.save(output_path)
print(f"Created: {output_path}")
print(f"Total slides: {len(prs.slides)}")
