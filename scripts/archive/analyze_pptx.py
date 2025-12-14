"""Extract detailed text content from PPTX for review."""
from pptx import Presentation
import json
import sys

def analyze_pptx(pptx_path):
    """Analyze PPTX content for final review."""
    prs = Presentation(pptx_path)
    
    analysis = {
        "total_slides": len(prs.slides),
        "slides": []
    }
    
    for slide_idx, slide in enumerate(prs.slides):
        slide_data = {
            "slide_index": slide_idx,
            "shapes": []
        }
        
        for shape_idx, shape in enumerate(slide.shapes):
            if not hasattr(shape, "text_frame") or not shape.has_text_frame:
                continue
            
            text_content = shape.text_frame.text
            para_count = len(shape.text_frame.paragraphs)
            
            shape_data = {
                "shape_index": shape_idx,
                "text": text_content,
                "char_count": len(text_content),
                "line_count": text_content.count('\n') + 1,
                "paragraph_count": para_count
            }
            
            slide_data["shapes"].append(shape_data)
        
        analysis["slides"].append(slide_data)
    
    return analysis

if __name__ == "__main__":
    pptx_file = sys.argv[1] if len(sys.argv) > 1 else "output_ppt/20251211_microsoft_intro_report.pptx"
    result = analyze_pptx(pptx_file)
    
    print(json.dumps(result, ensure_ascii=False, indent=2))
