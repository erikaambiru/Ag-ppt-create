# -*- coding: utf-8 -*-
# =============================================================================
# Ag-ppt-create - AI-powered PPTX generation pipeline
# https://github.com/aktsmm/Ag-ppt-create
# 
# Copyright (c) aktsmm. Licensed under CC BY-NC-SA 4.0.
# DO NOT MODIFY THIS HEADER BLOCK.
# =============================================================================
"""
Speaker Notes translator for PPTX files.
Extracts English speaker notes and applies Japanese translations.

Usage:
    # Extract notes to JSON
    python translate_notes.py extract input.pptx output_notes.json
    
    # Apply translated notes from JSON
    python translate_notes.py apply input.pptx notes_ja.json output.pptx
    
    # Quick translate with AI (requires manual edit of JSON)
    python translate_notes.py quick input.pptx output.pptx
"""
from pptx import Presentation
import json
import sys
import os


def extract_notes(input_file: str, output_json: str) -> dict:
    """
    Extract all speaker notes from PPTX to JSON file.
    
    Args:
        input_file: Path to input PPTX
        output_json: Path to output JSON file
        
    Returns:
        Dictionary with slide indices as keys and notes as values
    """
    prs = Presentation(input_file)
    notes = {}
    
    for i, slide in enumerate(prs.slides):
        if slide.has_notes_slide:
            notes_slide = slide.notes_slide
            text = notes_slide.notes_text_frame.text.strip()
            if text:
                notes[str(i)] = {
                    "original": text,
                    "translated": ""  # Placeholder for translation
                }
                print(f"Slide {i}: {len(text)} chars")
    
    with open(output_json, 'w', encoding='utf-8') as f:
        json.dump(notes, f, ensure_ascii=False, indent=2)
    
    print(f"\nExtracted {len(notes)} notes to: {output_json}")
    return notes


def apply_notes(input_file: str, notes_json: str, output_file: str) -> int:
    """
    Apply translated notes from JSON to PPTX.
    
    Args:
        input_file: Path to input PPTX
        notes_json: Path to JSON file with translations
        output_file: Path to output PPTX
        
    Returns:
        Number of notes applied
    """
    prs = Presentation(input_file)
    
    with open(notes_json, 'r', encoding='utf-8') as f:
        notes = json.load(f)
    
    applied_count = 0
    for slide_idx_str, note_data in notes.items():
        slide_idx = int(slide_idx_str)
        
        # Support both formats: {"translated": "..."} or just "..."
        if isinstance(note_data, dict):
            translated = note_data.get("translated", "")
        else:
            translated = note_data
        
        if not translated:
            continue
            
        if slide_idx < len(prs.slides):
            slide = prs.slides[slide_idx]
            # Ensure notes slide exists
            if not slide.has_notes_slide:
                slide.notes_slide  # This creates it
            notes_frame = slide.notes_slide.notes_text_frame
            notes_frame.text = translated
            applied_count += 1
            print(f"Applied notes to slide {slide_idx}")
    
    prs.save(output_file)
    print(f"\nApplied {applied_count} notes")
    print(f"Saved to: {output_file}")
    return applied_count


def translate_notes_legacy(input_file: str, output_file: str) -> None:
    """Legacy function: Translate speaker notes using hardcoded translations."""
    prs = Presentation(input_file)
    
    translated_count = 0
    for i, slide in enumerate(prs.slides):
        slide_num = i + 1
        if slide_num in NOTES_TRANSLATIONS:
            if slide.has_notes_slide:
                notes_slide = slide.notes_slide
                tf = notes_slide.notes_text_frame
                if tf:
                    tf.text = NOTES_TRANSLATIONS[slide_num]
                    translated_count += 1
                    print(f"Translated slide {slide_num}")
    
    prs.save(output_file)
    print(f"\nTranslated {translated_count} speaker notes")
    print(f"Saved to: {output_file}")


def print_usage():
    """Print usage information."""
    print("""
Speaker Notes Translator for PPTX

Usage:
    python translate_notes.py extract <input.pptx> <output.json>
        Extract speaker notes from PPTX to JSON for translation
        
    python translate_notes.py apply <input.pptx> <notes.json> <output.pptx>
        Apply translated notes from JSON to PPTX
        
    python translate_notes.py legacy <input.pptx> <output.pptx>
        Use hardcoded translations (BRK252 only)

JSON Format for apply:
    {
        "0": {"original": "English text", "translated": "日本語テキスト"},
        "1": "日本語テキスト"  // Simple format also supported
    }
""")


if __name__ == "__main__":
    if len(sys.argv) < 2:
        print_usage()
        sys.exit(1)
    
    command = sys.argv[1]
    
    if command == "extract":
        if len(sys.argv) < 4:
            print("Usage: python translate_notes.py extract <input.pptx> <output.json>")
            sys.exit(1)
        extract_notes(sys.argv[2], sys.argv[3])
        
    elif command == "apply":
        if len(sys.argv) < 5:
            print("Usage: python translate_notes.py apply <input.pptx> <notes.json> <output.pptx>")
            sys.exit(1)
        apply_notes(sys.argv[2], sys.argv[3], sys.argv[4])
        
    elif command == "legacy":
        if len(sys.argv) < 4:
            print("Usage: python translate_notes.py legacy <input.pptx> <output.pptx>")
            sys.exit(1)
        translate_notes_legacy(sys.argv[2], sys.argv[3])
        
    else:
        print(f"Unknown command: {command}")
        print_usage()
        sys.exit(1)


# Legacy support: hardcoded translations for BRK252
NOTES_TRANSLATIONS = {
    15: """これらは2025年9月の FabCon EU での最新の Purview + Fabric 発表です。
新しい「Public Preview」機能をデモします。""",

    16: """M365 と同様に、Fabric でも過剰共有の問題は現実です。
つまり、同様のプレイブックを適用できます：
- DSPM の Fabric リスク評価レポートで露出リスクを特定
- Information Protection で機密ラベルを定義
- または既存の M365 ラベルを再利用して Fabric アイテムに適用
- DLP で大規模にアクセスを制限""",

    17: """今後12〜18ヶ月で、80%のリーダーが増大する労働力需要に対応するためにエージェントを展開する予定です。
2028年までに、世界中の企業で13億以上のエージェントが使用されると予測しています。
この急速な採用は組織の運営方法を変革していますが、新たなリスクと要件も生み出しています。""",

    18: """データの保護 - AI の保護
このセクションでは、M365 Copilot がどのようにデータを使用し、どのように保護するかを説明します。
主な懸念事項：過剰共有とデータ損失""",

    19: """ほぼすべての組織で、データセキュリティは2つの補完的な役割が連携して機密情報を保護し、リスクに対応することに依存しています。
データセキュリティ管理者は、インシデントが発生する前に組織のデータセキュリティ態勢を形成し維持します。
この役割は、セキュリティポリシー（DLP、Insider Risk）の設計と調整、データ分類プログラムの推進、露出の削減に焦点を当てています。""",

    20: """ここで、DLP と IRM アラートのトリアージだけでなく、修復を推進する機能も追加したアラートトリアージエージェントの GA を発表します。
見てみましょう。""",

    21: """まず、データセキュリティアナリストとこれらのチームが直面している現実の課題、つまりアラート過負荷について話しましょう。
平均して、3件に1件以上のデータセキュリティアラートが完全に調査されていません。
経済面を見ると、アラートのトリアージにサードパーティベンダーを雇う平均コストは5万ドル、トリアージに必要な時間は21分です。
明らかに改善の機会があります。""",

    22: """データセキュリティ管理者をサポートするために、Data Security Posture Agent を発表できることを大変嬉しく思います。
このエージェントは、態勢のギャップを特定し、機密コンテンツを発見し、ポリシーの衛生状態を改善することで、データリスクを事前に管理するのに役立ちます。""",

    24: """可視性なくして保護なし。見えないものは保護できません。だからこそ、DSPM に AI 可観測性を導入しています。""",

    25: """データの保護 - AI の保護
M365 Copilot がどのようにデータを使用するか、そしてどのように保護するかの要約です。
主な懸念事項：過剰共有とデータ損失への対策""",
}
